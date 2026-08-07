import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Premium Color Palette (Same sleek cyber-dark theme)
    BG_COLOR = RGBColor(9, 13, 22)           # Slate dark background #090d16
    CARD_BG_COLOR = RGBColor(17, 24, 39)     # Panel card background #111827
    ACCENT_CYAN = RGBColor(6, 182, 212)      # Cyber Cyan #06b6d4
    ACCENT_INDIGO = RGBColor(99, 102, 241)   # Neon Indigo #6366f1
    TEXT_PRIMARY = RGBColor(248, 250, 252)    # Clean White #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)      # Soft Gray #94a3b8
    ACCENT_EMERALD = RGBColor(16, 185, 129)  # Emerald Green
    ACCENT_AMBER = RGBColor(245, 158, 11)     # Warning Amber
    ACCENT_PURPLE = RGBColor(139, 92, 246)    # Blocked Purple
    
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_bg_and_header(slide, title_text, category_text="AE-03 ORCHESTRATOR"):
        # Dark Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Sleek Neon Header Accent Bar (Top Border)
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ACCENT_CYAN
        top_bar.line.fill.background()
        
        # Category Tracker Label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = "Consolas"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = ACCENT_INDIGO
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Segoe UI"
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_PRIMARY

    # ====================================================
    # SLIDE 1: Title Slide (High Impact Cover)
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_CYAN
    stripe.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "AE-03 ORCHESTRATOR"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "Unified Agent Form Orchestration Engine"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_INDIGO
    p2.space_before = Pt(6)
    
    meta_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11.0), Inches(2.5))
    meta_tf = meta_box.text_frame
    meta_tf.word_wrap = True
    
    fields = [
        ("Team Name", "[Enter your Team Name]"),
        ("Team Leader Name", "[Enter Team Leader Name]"),
        ("Problem Statement", "Traditional multi-agent pipelines lack strict tool sandboxing, input/output validation, budget enforcement, and human-in-the-loop audit controls. AE-03 delivers a provider-abstracted engine that compiles goals into safe, typed execution graphs.")
    ]
    for i, (label, val) in enumerate(fields):
        p = meta_tf.add_paragraph() if i > 0 else meta_tf.paragraphs[0]
        p.text = f"{label}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        
        # Add value text as unbolded
        run = p.add_run()
        run.text = val
        run.font.name = "Segoe UI"
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # ====================================================
    # SLIDE 2: Team Members
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Orchestration Team", "MEMBERS REGISTRY")
    
    # Render cards for 4 team members
    card_width = Inches(2.7)
    card_height = Inches(4.5)
    gap = Inches(0.2)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)
    
    members = [
        {"role": "TEAM LEADER", "name": "[Leader Name]", "college": "[College Name]", "linkedin": "[LinkedIn link]"},
        {"role": "DEVELOPER 1", "name": "[Member Name]", "college": "[College Name]", "linkedin": "[LinkedIn link]"},
        {"role": "DEVELOPER 2", "name": "[Member Name]", "college": "[College Name]", "linkedin": "[LinkedIn link]"},
        {"role": "DEVELOPER 3", "name": "[Member Name]", "college": "[College Name]", "linkedin": "[LinkedIn link]"}
    ]
    
    for i, m in enumerate(members):
        left_pos = start_left + i * (card_width + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, card_width, card_height)
        card.fill.solid()
        if i == 0:
            card.fill.fore_color.rgb = CARD_BG_COLOR
            card.line.color.rgb = ACCENT_CYAN
            card.line.width = Pt(1.5)
        else:
            card.fill.fore_color.rgb = CARD_BG_COLOR
            card.line.color.rgb = RGBColor(51, 65, 85)
            
        tf = card.text_frame
        tf.word_wrap = True
        
        # Role header
        p = tf.paragraphs[0]
        p.text = m["role"]
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if i == 0 else ACCENT_INDIGO
        p.space_after = Pt(18)
        
        # Name
        p2 = tf.add_paragraph()
        p2.text = m["name"]
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_PRIMARY
        p2.space_after = Pt(14)
        
        # College
        p3 = tf.add_paragraph()
        p3.text = "College:"
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(10)
        p3.font.color.rgb = TEXT_MUTED
        
        p3_val = tf.add_paragraph()
        p3_val.text = m["college"]
        p3_val.font.name = "Segoe UI"
        p3_val.font.size = Pt(12)
        p3_val.font.bold = True
        p3_val.font.color.rgb = TEXT_PRIMARY
        p3_val.space_after = Pt(14)
        
        # Linkedin
        p4 = tf.add_paragraph()
        p4.text = "LinkedIn:"
        p4.font.name = "Segoe UI"
        p4.font.size = Pt(10)
        p4.font.color.rgb = TEXT_MUTED
        
        p4_val = tf.add_paragraph()
        p4_val.text = m["linkedin"]
        p4_val.font.name = "Segoe UI"
        p4_val.font.size = Pt(11)
        p4_val.font.color.rgb = ACCENT_CYAN
        p4_val.font.underline = True

    # ====================================================
    # SLIDE 3: Brief about the Idea
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Unified Agent Orchestrator Concept", "CORE METHODOLOGY")
    
    # Large left card: Core Concept
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG_COLOR
    left_card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = left_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "THE CONCEPT"
    p.font.name = "Consolas"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "AE-03 is a provider-abstracted multi-agent orchestration console designed to bridge the trust gap in AI autonomy."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_PRIMARY
    p2.space_after = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Users input goals in plain language. The engine parses them, outputs structured execution flows (nodes and edges), and executes them. Every data handoff is validated against JSON schemas at the message bus, preventing code corruption or injection escapes."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = TEXT_MUTED
    p3.line_spacing = 1.3
    
    # Right column: 3 Key Modules
    modules = [
        ("01", "Goal Compiler", "Translates natural language text goals into human-approvable execution graph schema structures."),
        ("02", "Deterministic Engine", "Manages sequential, parallel, retry loops, and human-in-the-loop review checks with SQLite states."),
        ("03", "Security Sandbox", "Isolates keys, filters agent tools by scopes, and blocks adversarial overrides at the message bus.")
    ]
    
    for i, (num, title, desc) in enumerate(modules):
        m_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8 + i * 1.65), Inches(5.7), Inches(1.5))
        m_card.fill.solid()
        m_card.fill.fore_color.rgb = CARD_BG_COLOR
        m_card.line.color.rgb = RGBColor(51, 65, 85)
        
        m_tf = m_card.text_frame
        m_tf.word_wrap = True
        
        # Grid layout inside card using tabs or spacing
        p = m_tf.paragraphs[0]
        p.text = f"{num} | {title.upper()}"
        p.font.name = "Consolas"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO
        p.space_after = Pt(6)
        
        p2 = m_tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.line_spacing = 1.2

    # ====================================================
    # SLIDE 4: Opportunity & USP
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Security-by-Design & Unique Value", "MARKET ADVANTAGE")
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG_COLOR
    left_card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = left_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HOW WE SOLVE AUTONOMY RISKS"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(14)
    
    solves = [
        "Typed Handoff Contracts: Stops data format drift by enforcing strict JSON input/output schemas on every agent connection.",
        "Tool Allowlisting: Intercepts and blocks unauthorized tool calls (e.g. system CLI escapes) at the message bus layer.",
        "Token Cost Ledger: Enforces pre-declared budget limits on agents, preventing infinite run loops and cost overruns.",
        "Sandboxed Directory: Boundaries file reads/writes strictly within a declared local agent workspace folder."
    ]
    for s in solves:
        p = tf.add_paragraph()
        p.text = chr(9656) + " " + s
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG_COLOR
    right_card.line.color.rgb = ACCENT_CYAN
    
    tf2 = right_card.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "THE UNIQUE SELLING PROPOSITION (USP)"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)
    
    usps = [
        "Goal-to-Graph Propose & Lock: Compiles user intent into a visual plan that a human must review and approve before any execution.",
        "Deterministic Replay & Stats: Saves execution events to SQLite, allowing logs to be replayed to check latency and costs.",
        "Adversarial Interception Demo: Runs live tests showing the engine trapping and displaying out-of-scope command breaches.",
        "Marginal Agent Value Evaluation: Evaluates single vs. multi-agent configurations, reporting when extra agents are not worth it."
    ]
    for u in usps:
        p = tf2.add_paragraph()
        p.text = chr(9670) + " " + u
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)
        p.line_spacing = 1.2

    # ====================================================
    # SLIDE 5: Process Flow
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Deterministic Control Pipeline Flow", "ENGINE DATAFLOW")
    
    flow_steps = [
        ("Goal Input", "User states goal in natural language"),
        ("Compile Plan", "Engine translates intent into JSON nodes/edges"),
        ("Approve Gate", "Human audits budgets, permissions, and locks schema"),
        ("Step Dispatch", "Engine schedules parallel/sequential agent tasks"),
        ("Safety Verify", "Secret broker masks keys; tool allowlist checked"),
        ("Handoff Check", "Data output validated against target input schema"),
        ("Audit / Log", "SQLite saves event traces and registers artifacts")
    ]
    
    box_width = Inches(1.5)
    box_height = Inches(3.2)
    gap = Inches(0.16)
    start_left = Inches(0.8)
    top_pos = Inches(2.2)
    
    for i, (title, desc) in enumerate(flow_steps):
        left_pos = start_left + i * (box_width + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG_COLOR
        box.line.color.rgb = ACCENT_CYAN if i in [2, 4, 5] else RGBColor(51, 65, 85)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"STEP 0{i+1}"
        p.font.name = "Consolas"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if i % 2 == 0 else ACCENT_INDIGO
        p.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_PRIMARY
        p2.space_after = Pt(8)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.line_spacing = 1.2
        
        # Add visual arrow except for last step
        if i < len(flow_steps) - 1:
            arrow = slide.shapes.add_textbox(left_pos + box_width, top_pos + Inches(1.2), gap, Inches(0.8))
            arrow_tf = arrow.text_frame
            arrow_p = arrow_tf.paragraphs[0]
            arrow_p.text = chr(8594)
            arrow_p.alignment = PP_ALIGN.CENTER
            arrow_p.font.name = "Segoe UI"
            arrow_p.font.size = Pt(18)
            arrow_p.font.bold = True
            arrow_p.font.color.rgb = ACCENT_CYAN

    # ====================================================
    # SLIDE 6: UI Wireframes & Layout Design
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Control Console Wireframe Mockups", "DASHBOARD INTERFACE")
    
    # Large card: UI Screens and Layouts
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG_COLOR
    left_card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = left_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DESIGN AESTHETICS"
    p.font.name = "Consolas"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "Premium, information-dense cybernetic control room console theme."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_PRIMARY
    p2.space_after = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Using a deep slate-dark (#090d16) canvas with clean glow cards and neon borders indicating status (Green=Success, Amber=Retry, Red=Failure/Review, Purple=Blocked). The layout avoids unnecessary chatbot conversation bubbles and focuses on pipeline structure, cost metrics, and log events."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = TEXT_MUTED
    p3.line_spacing = 1.3
    
    # Right column: 3 Screens
    screens = [
        ("01", "COMPILER & APPROVAL", "Goal input, compiled node list, templates, budgets, and the Lock & Approve gate button."),
        ("02", "EXECUTION BOARD", "Visual SVG graph nodes, active status pulses, live token ledgers, events console, and review forms."),
        ("03", "TRACE & REPLAY CENTER", "Picker showing historical run runs, detailed logs side-by-side, and deterministic replay comparisons.")
    ]
    
    for i, (num, title, desc) in enumerate(screens):
        s_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8 + i * 1.65), Inches(5.7), Inches(1.5))
        s_card.fill.solid()
        s_card.fill.fore_color.rgb = CARD_BG_COLOR
        s_card.line.color.rgb = ACCENT_CYAN if i == 1 else RGBColor(51, 65, 85)
        
        s_tf = s_card.text_frame
        s_tf.word_wrap = True
        
        p = s_tf.paragraphs[0]
        p.text = f"{num} | {title}"
        p.font.name = "Consolas"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if i == 1 else TEXT_MUTED
        p.space_after = Pt(6)
        
        p2 = s_tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.line_spacing = 1.2

    # ====================================================
    # SLIDE 7: Technical Architecture
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Technical Architecture & Module Topology", "SYSTEM TOPOLOGY")
    
    # Left Card: Backend Core Modules
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG_COLOR
    left_card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = left_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BACKEND ENGINE LAYERS (FASTAPI)"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(14)
    
    layers = [
        "Compiler Module: Evaluates user prompts to output standard JSON graphs.",
        "Execution Machine: Step-by-step runner handling parallel, sequential, and retries.",
        "Security boundary: Checks tool names, sanitizes payload injection overrides.",
        "Secret Broker: Keeps API keys, maps calls with scope credentials.",
        "SQLite Persistence: Stores agent profiles, runs, events, and artifacts."
    ]
    for l in layers:
        p = tf.add_paragraph()
        p.text = chr(9656) + " " + l
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        
    # Right Card: Frontend UI & Channels
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG_COLOR
    right_card.line.color.rgb = ACCENT_CYAN
    
    tf2 = right_card.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "FRONTEND & MESSAGE channels (REACT)"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)
    
    fe_layers = [
        "Goal Editor: Lets users inspect and modify graph variables in real time.",
        "Graph Visualiser: Custom render connecting nodes and animating running paths.",
        "Audit Log Stream: Streams runtime JSON events (starts, retries, blocks).",
        "Artifact Previews: Triggers popups verifying payload validation stamps.",
        "Evaluation View: Harness comparing single vs multi-agent execution."
    ]
    for f in fe_layers:
        p = tf2.add_paragraph()
        p.text = chr(9670) + " " + f
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)
        p.line_spacing = 1.2

    # ====================================================
    # SLIDE 8: Technologies Used
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Development Technologies Stack", "TECH STACK")
    
    techs = [
        ("FastAPI", "Python", "High-performance API server with automatic OpenAPI docs and routing dependencies."),
        ("SQLite", "Zero-Ops DB", "Transactional SQL storage storing graphs, runs, events logs, and artifacts."),
        ("React + TS", "Frontend", "Single-Page Application dashboard with Vite compiler and state hooks."),
        ("JSON Schema", "Validation", "Libraries (`jsonschema` in Python, schemas in React) verifying data contracts."),
        ("Ollama / APIs", "LLM Providers", "Local Llama integration with API fallbacks (OpenAI, Anthropic) and Mock Resiliency."),
        ("Lucide / CSS", "Aesthetics", "Lucide icons, custom CSS variables, and animated cybernetic elements.")
    ]
    
    for i, (name, role, desc) in enumerate(techs):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + col * Inches(3.9)
        top = Inches(1.8) + row * Inches(2.4)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG_COLOR
        card.line.color.rgb = RGBColor(51, 65, 85)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        
        p = ctf.paragraphs[0]
        p.text = name
        p.font.name = "Segoe UI"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        
        p_role = ctf.add_paragraph()
        p_role.text = role.upper()
        p_role.font.name = "Consolas"
        p_role.font.size = Pt(10.5)
        p_role.font.bold = True
        p_role.font.color.rgb = ACCENT_CYAN
        p_role.space_after = Pt(8)
        
        p_desc = ctf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.line_spacing = 1.2

    # ====================================================
    # SLIDE 9: Cost & Performance Harness Analysis
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg_and_header(slide, "Cost & Performance Harness Evaluation", "BENCHMARKS & AUDITING")
    
    # Left Card: Cost Analysis
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG_COLOR
    left_card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = left_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "RUN-TIME COST LEDGER"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)
    
    costs = [
        "SQLite Database Hosting: $0 (zero-ops local file).",
        "Ollama Local LLM Runs: $0 (local hardware execution).",
        "OpenAI API cost check: GPT-4o-mini is ~$0.15 / 1M input tokens. Average compiler call is less than $0.0001.",
        "Anthropic API cost check: Claude-3-5-sonnet is ~$3.0 / 1M input tokens. Average review verification step is ~$0.002.",
        "Total Sandbox Cost: $0 (unlimited local testing using the robust rule-based mock provider fallback)."
    ]
    for c in costs:
        p = tf.add_paragraph()
        p.text = chr(9656) + " " + c
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        
    # Right Card: Performance Harness
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG_COLOR
    right_card.line.color.rgb = ACCENT_INDIGO
    
    tf2 = right_card.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "EVALUATION HARNESS AUDIT"
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(14)
    
    evals = [
        "Single vs. Multi-Agent Benchmark: Evaluates execution across 3 default task templates, logging cost, latency and success rates.",
        "Marginal Agent Value reporting: Reports when extra agents add complexity without improving accuracy (preventing bloat).",
        "Recovery Rate audit: Logs recovery frequency when parallel paths hit timeouts, demonstrating retry loop integrity.",
        "Deterministic replay tests: Compares original vs replayed logs to detect logic drift in model configurations."
    ]
    for e in evals:
        p = tf2.add_paragraph()
        p.text = chr(9670) + " " + e
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)
        p.line_spacing = 1.2

    # ====================================================
    # SLIDE 10: Thank You
    # ====================================================
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_CYAN
    stripe.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "THANK YOU"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(72)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "AE-03: Unified Agent Form Orchestrator"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_INDIGO
    p2.space_before = Pt(6)
    
    meta_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11.0), Inches(2.0))
    meta_tf = meta_box.text_frame
    meta_tf.word_wrap = True
    
    p = meta_tf.paragraphs[0]
    p.text = "Contact Email: [Enter your Email]\n"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    run = p.add_run()
    run.text = "GitHub Repository: https://github.com/Goldengrab/YCCE\n"
    run.font.name = "Segoe UI"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TEXT_PRIMARY
    
    run2 = p.add_run()
    run2.text = "Advanced Agentic Systems Challenge — AE-03 Solution Blueprint"
    run2.font.name = "Segoe UI"
    run2.font.size = Pt(13)
    run2.font.italic = True
    run2.font.color.rgb = TEXT_MUTED
    
    prs.save("AE03_Orchestrator_Pitch_Deck.pptx")
    print("Presentation created successfully as 'AE03_Orchestrator_Pitch_Deck.pptx'!")

if __name__ == "__main__":
    create_deck()
